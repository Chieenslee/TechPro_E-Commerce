import collections
import collections.abc
import pptx
import os

prs = pptx.Presentation(r"d:\My\CNPM\TechPro_E-Commerce\Doc\BCTD2.pptx")

with open(r"d:\My\CNPM\TechPro_E-Commerce\pptx_content_bctd2_utf8.txt", "w", encoding="utf-8") as f:
    for i, slide in enumerate(prs.slides):
        f.write(f"--- Slide {i+1} ---\n")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                f.write(shape.text.strip() + "\n")
            elif shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
                f.write("[PICTURE]\n")
            elif shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP:
                f.write("[GROUP OF SHAPES/DIAGRAM]\n")
            elif shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.CHART:
                f.write("[CHART]\n")
            elif shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:
                f.write("[TABLE]\n")
                for row in shape.table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip().replace('\n', ' '))
                    f.write(" | ".join(row_data) + "\n")
        f.write("\n")
