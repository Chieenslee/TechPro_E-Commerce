import collections
import collections.abc
import pptx
import os

prs = pptx.Presentation(r"d:\My\CNPM\TechPro_E-Commerce\Doc\Nhom5_BCCD1 (1).pptx")
output_dir = r"d:\My\CNPM\TechPro_E-Commerce\Doc\slides"
os.makedirs(output_dir, exist_ok=True)

for i, slide in enumerate(prs.slides):
    for j, shape in enumerate(slide.shapes):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            image = shape.image
            image_bytes = image.blob
            image_ext = image.ext
            image_filename = os.path.join(output_dir, f"slide_{i+1}.{image_ext}")
            with open(image_filename, "wb") as f:
                f.write(image_bytes)
            print(f"Saved {image_filename}")
