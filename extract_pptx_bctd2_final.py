import collections
import collections.abc
import pptx

prs = pptx.Presentation(r"d:\My\CNPM\TechPro_E-Commerce\Doc\BCTD2_Final.pptx")
with open(r"d:\My\CNPM\TechPro_E-Commerce\pptx_content_bctd2_final_utf8.txt", "w", encoding="utf-8") as f:
    for i, slide in enumerate(prs.slides):
        f.write(f"--- Slide {i+1} ---\n")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                f.write(shape.text.strip() + "\n")
            if shape.has_table:
                f.write("[TABLE]\n")
                for row in shape.table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip().replace('\n', ' '))
                    f.write(" | ".join(row_data) + "\n")
            if shape.shape_type == 6:  # GROUP
                for sub_shape in shape.shapes:
                    if hasattr(sub_shape, "text") and sub_shape.text.strip():
                        f.write(sub_shape.text.strip() + "\n")
                    if sub_shape.has_table:
                        f.write("[TABLE]\n")
                        for row in sub_shape.table.rows:
                            row_data = []
                            for cell in row.cells:
                                row_data.append(cell.text.strip().replace('\n', ' '))
                            f.write(" | ".join(row_data) + "\n")
        f.write("\n")
