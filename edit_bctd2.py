import collections
import collections.abc
from pptx import Presentation
import os

prs = Presentation(r"d:\My\CNPM\TechPro_E-Commerce\Doc\BCTD2.pptx")

def replace_text(shape, replacements):
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            for old_text, new_text in replacements.items():
                if old_text in run.text:
                    run.text = run.text.replace(old_text, new_text)

def process_table(shape, replacements):
    if not shape.has_table:
        return
    for row in shape.table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    for old_text, new_text in replacements.items():
                        if old_text in run.text:
                            run.text = run.text.replace(old_text, new_text)

replacements = {
    "mock_token_123": "JWT Token",
    "mock_token": "JWT Token",
    "Chưa có: xUnit/NUnit": "Đã có Unit Test bằng xUnit",
    "Mock auth · chưa JWT/OAuth": "(Đã khắc phục) Đã tích hợp JWT Auth",
    "Chưa unit/E2E tests trong repo": "(Đã khắc phục) Đã bổ sung xUnit tests",
    "Thay mock_token_123 bằng JWT + refresh token": "Cấu hình thời gian hết hạn cho JWT Token"
}

for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            replace_text(shape, replacements)
        if shape.has_table:
            process_table(shape, replacements)
        if shape.shape_type == 6:  # GROUP
            for sub_shape in shape.shapes:
                if sub_shape.has_text_frame:
                    replace_text(sub_shape, replacements)
                if sub_shape.has_table:
                    process_table(sub_shape, replacements)

prs.save(r"d:\My\CNPM\TechPro_E-Commerce\Doc\BCTD2_Final.pptx")
print("Done saving improved BCTD2_Final.pptx")
